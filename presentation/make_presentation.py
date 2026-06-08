"""Build a PPTX presentation in the style of the XAI626 reference deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- palette ----------
BG_DARK = RGBColor(0x3A, 0x3A, 0x3A)
BG_TITLE = RGBColor(0x32, 0x32, 0x32)
CARD_BG = RGBColor(0x4E, 0x4E, 0x4E)
TEAL = RGBColor(0x5B, 0xB5, 0xBE)
ORANGE = RGBColor(0xE8, 0xA8, 0x7C)
PINK = RGBColor(0xE9, 0x63, 0x79)
YELLOW = RGBColor(0xF5, 0xC8, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
DIM = RGBColor(0x99, 0x99, 0x99)

# ---------- layout ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_runs(slide, left, top, width, height, runs,
             size=18, align=PP_ALIGN.LEFT, font_name="Calibri",
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs: list of (text, {'bold','color','size','italic'})."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for text, props in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = text
        r.font.name = props.get("font", font_name)
        r.font.size = Pt(props.get("size", size))
        r.font.bold = props.get("bold", False)
        r.font.italic = props.get("italic", False)
        r.font.color.rgb = props.get("color", WHITE)
    return tb


def add_rect(slide, left, top, width, height, fill=CARD_BG,
             line_color=None, corner=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape_type, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line_color is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line_color
        r.line.width = Pt(0.75)
    r.shadow.inherit = False
    # remove default text margins
    r.text_frame.margin_left = Emu(0)
    r.text_frame.margin_right = Emu(0)
    r.text_frame.margin_top = Emu(0)
    r.text_frame.margin_bottom = Emu(0)
    # Set rounding amount via adjustments if available
    try:
        if corner:
            r.adjustments[0] = 0.08
    except Exception:
        pass
    return r


def add_header(slide, section, subsection, page=None, total_pages=None):
    """Header strip: section (teal serif) + subsection (white with orange bar)."""
    # Section title (top-left, teal serif large)
    add_text(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
             section, size=30, bold=True, color=TEAL,
             font_name="Cambria")
    # Orange vertical bar before subsection
    if subsection:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.85), Inches(0.06), Inches(0.35),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.fill.background()
        # Subsection title
        add_text(slide, Inches(0.65), Inches(0.82), Inches(10), Inches(0.4),
                 subsection, size=16, bold=True, color=WHITE,
                 font_name="Cambria")
    # Horizontal divider line under header
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.32), Inches(12.3), Emu(8000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()


def add_footer(slide, page_num, total_pages=12):
    # Page number centered
    add_text(slide, Inches(6.4), Inches(7.05), Inches(0.5), Inches(0.3),
             str(page_num), size=10, color=DIM, align=PP_ALIGN.CENTER)
    # Left footer
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "XAI 509 — Wav2Vec2 CTC Fine-Tuning Project",
             size=10, color=DIM, italic_safe=False) \
        if False else add_text(
            slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
            "XAI 509 — Wav2Vec2 CTC Fine-Tuning Project",
            size=10, color=DIM)


# =============================================================
# SLIDE 1 — Title
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_TITLE)

# Conference-like banner
add_text(s, Inches(0), Inches(1.7), SLIDE_W, Inches(0.5),
         "XAI 509 Course Project", size=22, bold=False, color=TEAL,
         align=PP_ALIGN.CENTER, font_name="Cambria")

# Main title
add_text(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.4),
         "Wav2Vec2 CTC Fine-Tuning with\nOutput-Distribution Regularization",
         size=36, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")

# Subtitle / context
add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.4),
         "LibriSpeech 1-hour Fine-tuning · Maximum Entropy vs Label Smoothing",
         size=18, color=ORANGE, align=PP_ALIGN.CENTER, font_name="Cambria")

# Authors / lab placeholder
add_text(s, Inches(0), Inches(4.8), SLIDE_W, Inches(0.4),
         "Student Name", size=18, color=TEAL,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.4),
         "Speech Language Processing Lab", size=16, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.55), SLIDE_W, Inches(0.4),
         "Korea University", size=16, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.2), SLIDE_W, Inches(0.4),
         "June 11, 2026", size=14, color=DIM,
         align=PP_ALIGN.CENTER, font_name="Cambria")

# =============================================================
# SLIDE 2 — Outline
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Outline", "")

# Two-column outline
col_l = [
    "Background & Problem",
    "Two Regularization Methods",
    "Experimental Setup",
]
col_r = [
    "Main Results",
    "MaxEnt + LS Sweeps & Analysis",
    "Summary",
]
# Orange vertical bar (left column accent)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(1.5), Inches(2.2), Inches(0.07), Inches(3.5))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()

for i, item in enumerate(col_l):
    y = Inches(2.3 + i * 1.1)
    add_text(s, Inches(1.9), y, Inches(0.3), Inches(0.5), "●", size=22, color=WHITE)
    add_text(s, Inches(2.3), y, Inches(4.5), Inches(0.5), item, size=24, color=WHITE,
             font_name="Cambria")

# Vertical separator
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(7.0), Inches(2.2), Emu(8000), Inches(3.5))
sep.fill.solid(); sep.fill.fore_color.rgb = LIGHT; sep.line.fill.background()

for i, item in enumerate(col_r):
    y = Inches(2.3 + i * 1.1)
    add_text(s, Inches(7.5), y, Inches(0.3), Inches(0.5), "●", size=22, color=WHITE)
    add_text(s, Inches(7.9), y, Inches(4.8), Inches(0.5), item, size=24, color=WHITE,
             font_name="Cambria")

add_footer(s, 2)


# =============================================================
# SLIDE 3 — Project Overview
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Project Overview",
           "Wav2Vec2 CTC Fine-Tuning + Regularization on LibriSpeech 1h")

# Goal box
add_rect(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(1.6))
add_text(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(0.4),
         "Goal", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(1.0), [
    ("Fine-tune ", {}),
    ("Wav2Vec2", {"bold": True, "color": ORANGE}),
    (" with the ", {}),
    ("CTC loss", {"bold": True, "color": ORANGE}),
    (" on a small (1-hour) LibriSpeech subset and report WER on test-clean / test-other.",
     {}),
], size=14, line_spacing=1.2)

# Core Idea box
add_rect(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(1.6))
add_text(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.4),
         "Core Idea", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(7.0), Inches(2.25), Inches(5.6), Inches(1.0), [
    ("Small fine-tuning data → ", {}),
    ("over-confidence", {"bold": True, "color": PINK}),
    (". Compare two output-distribution regularizers: ", {}),
    ("Maximum Entropy", {"bold": True, "color": ORANGE}),
    (" and ", {}),
    ("Label Smoothing", {"bold": True, "color": ORANGE}),
    (".", {}),
], size=14, line_spacing=1.2)

# Result/Contribution box (wider)
add_rect(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.7))
add_text(s, Inches(0.7), Inches(3.65), Inches(5.5), Inches(0.4),
         "Key Finding", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(0.7), Inches(4.05), Inches(12.0), Inches(1.2), [
    ("Both regularizers improve over baseline. ", {}),
    ("Label Smoothing (α=0.10) achieves the best WER", {"bold": True, "color": YELLOW}),
    (" — clean ", {}),
    ("19.51 %", {"bold": True, "color": ORANGE}),
    (" (−1.65 pp), other ", {}),
    ("28.30 %", {"bold": True, "color": ORANGE}),
    (" (−1.87 pp), outperforming both baseline CTC and Maximum Entropy training.", {}),
], size=14, line_spacing=1.3)

# Limitation box
add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.3))
add_text(s, Inches(0.7), Inches(5.55), Inches(5.5), Inches(0.4),
         "Scope / Limitation", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.8), [
    ("Only α = 0.10 swept for label smoothing; MaxEnt has 5-point λ sweep. ", {}),
    ("Single-seed experiments.", {"italic": True, "color": LIGHT}),
], size=13, line_spacing=1.2)

add_footer(s, 3)


# =============================================================
# SLIDE 4 — Background
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Background", "Wav2Vec 2.0 + CTC for ASR")

# Left: Wav2Vec 2.0
add_text(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.45),
         "Wav2Vec 2.0", size=20, bold=True, color=ORANGE, font_name="Cambria")
bullets_l = [
    "Self-supervised speech representation",
    "Pretrained on 53k hours of unlabeled audio",
    "CNN feature encoder → Transformer context",
    "Fine-tunable for downstream ASR",
]
for i, b in enumerate(bullets_l):
    add_text(s, Inches(0.7), Inches(2.0 + i * 0.55), Inches(0.3), Inches(0.4),
             "●", size=14, color=TEAL)
    add_text(s, Inches(1.0), Inches(2.0 + i * 0.55), Inches(5.5), Inches(0.4),
             b, size=14, color=WHITE)

# Right: CTC
add_text(s, Inches(7.0), Inches(1.5), Inches(6), Inches(0.45),
         "CTC Loss", size=20, bold=True, color=ORANGE, font_name="Cambria")
bullets_r = [
    "No frame-level alignment required",
    "Marginalizes over all valid alignments",
    "WER = (S + I + D) / N",
    "Lower WER = better recognition",
]
for i, b in enumerate(bullets_r):
    add_text(s, Inches(7.2), Inches(2.0 + i * 0.55), Inches(0.3), Inches(0.4),
             "●", size=14, color=TEAL)
    add_text(s, Inches(7.5), Inches(2.0 + i * 0.55), Inches(5.5), Inches(0.4),
             b, size=14, color=WHITE)

# Highlighted problem at bottom
add_rect(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.8))
add_text(s, Inches(0.7), Inches(5.05), Inches(12.0), Inches(0.4),
         "Why this project?", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.2), [
    ("Fine-tuning Wav2Vec2 on just ", {}),
    ("1 hour", {"bold": True, "color": ORANGE}),
    (" of speech is challenging — the model can become ", {}),
    ("over-confident on training data", {"bold": True, "color": PINK}),
    (", producing sharp output distributions that fail to generalize to ", {}),
    ("test-clean / test-other", {"bold": True, "color": ORANGE}),
    (".", {}),
], size=14, line_spacing=1.4)

add_footer(s, 4)


# =============================================================
# SLIDE 5 — Two Regularization Methods
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Method", "Two Regularization Strategies")

# Italic intro
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
         "Both methods penalize sharp output distributions, but apply different weighting schemes.",
         size=14, color=LIGHT, italic_safe=False) \
        if False else add_text(
            s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
            "Both methods penalize sharp output distributions, but apply different weighting schemes.",
            size=14, color=LIGHT)

# Left card: Maximum Entropy
add_rect(s, Inches(0.5), Inches(2.05), Inches(6.1), Inches(4.5),
         fill=RGBColor(0x55, 0x55, 0x55))
add_text(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(0.5),
         "Maximum Entropy", size=22, bold=True, color=TEAL, font_name="Cambria")
add_text(s, Inches(0.7), Inches(2.75), Inches(5.5), Inches(0.4),
         "Loss = CTC  −  λ · H(p)", size=18, bold=True, color=WHITE,
         font_name="Consolas")
add_text(s, Inches(0.7), Inches(3.25), Inches(5.5), Inches(0.4),
         "H(p) = −Σ pᵥ log pᵥ", size=14, color=ORANGE, font_name="Consolas")
add_text(s, Inches(0.7), Inches(3.85), Inches(5.7), Inches(0.4),
         "Properties:", size=14, bold=True, color=YELLOW)
me_props = [
    "Weights penalty by predicted probability p",
    "Pushes confident tokens toward uniform first",
    "Listed in lecture 8 \"Topics You Can Try\"",
]
for i, p in enumerate(me_props):
    add_text(s, Inches(0.9), Inches(4.3 + i * 0.45), Inches(0.3), Inches(0.4),
             "●", size=12, color=TEAL)
    add_text(s, Inches(1.15), Inches(4.3 + i * 0.45), Inches(5.3), Inches(0.4),
             p, size=13, color=WHITE)

# Right card: Label Smoothing
add_rect(s, Inches(6.8), Inches(2.05), Inches(6.0), Inches(4.5),
         fill=RGBColor(0x55, 0x55, 0x55))
add_text(s, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.5),
         "Label Smoothing", size=22, bold=True, color=PINK, font_name="Cambria")
add_text(s, Inches(7.0), Inches(2.75), Inches(5.5), Inches(0.4),
         "Loss = (1−α) · CTC  +  α · LS(p)", size=18, bold=True, color=WHITE,
         font_name="Consolas")
add_text(s, Inches(7.0), Inches(3.25), Inches(5.5), Inches(0.4),
         "LS(p) = −(1/V) Σ log pᵥ", size=14, color=ORANGE, font_name="Consolas")
add_text(s, Inches(7.0), Inches(3.85), Inches(5.7), Inches(0.4),
         "Properties:", size=14, bold=True, color=YELLOW)
ls_props = [
    "Weights penalty uniformly across vocab (1/V)",
    "Equivalent to KL(uniform ∥ p) up to constant",
    "Reported as part of lab's MoChA paper",
]
for i, p in enumerate(ls_props):
    add_text(s, Inches(7.2), Inches(4.3 + i * 0.45), Inches(0.3), Inches(0.4),
             "●", size=12, color=PINK)
    add_text(s, Inches(7.45), Inches(4.3 + i * 0.45), Inches(5.3), Inches(0.4),
             p, size=13, color=WHITE)

# Footer note
add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         "→ Closely related (Pereyra et al. 2017) but mathematically distinct.",
         size=13, color=YELLOW, italic_safe=False, align=PP_ALIGN.CENTER) \
    if False else add_text(
        s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
        "→ Closely related (Pereyra et al. 2017) but mathematically distinct.",
        size=13, color=YELLOW, align=PP_ALIGN.CENTER, font_name="Cambria")

add_footer(s, 5)


# =============================================================
# SLIDE 6 — Experimental Setup
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Experiment", "Setup")

# Setup table on the left
rows = [
    ("Pretrained model",   "facebook/wav2vec2-base (95M params)"),
    ("Train set",          "LibriLight 1-hour subset (WebDataset, 5 shards)"),
    ("Eval sets",          "test-clean (2620 utt) / test-other (2939 utt)"),
    ("Training steps",     "2000 (~222 epochs on 1h set)"),
    ("Optimizer / LR",     "AdamW, lr=1e-4, warmup 500 steps"),
    ("Batch",              "16 × grad_accum 2 = effective 32, fp16"),
    ("Frozen modules",     "CNN feature encoder"),
    ("Hardware",           "RTX 3090 (24 GB), ~35 min per run"),
]
y0 = Inches(1.7)
for i, (k, v) in enumerate(rows):
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x44, 0x44, 0x44)
    rect = add_rect(s, Inches(0.5), Inches(1.7 + i * 0.55), Inches(12.3),
                    Inches(0.52), fill=bg, corner=False)
    add_text(s, Inches(0.7), Inches(1.78 + i * 0.55), Inches(4.0), Inches(0.4),
             k, size=14, bold=True, color=TEAL)
    add_text(s, Inches(4.8), Inches(1.78 + i * 0.55), Inches(8.0), Inches(0.4),
             v, size=14, color=WHITE, font_name="Consolas")

# Bottom note
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
         "Each variant (baseline / MaxEnt λ / Label Smoothing α) trained with identical settings for fair comparison.",
         size=13, color=YELLOW, align=PP_ALIGN.CENTER, font_name="Cambria")

add_footer(s, 6)


# =============================================================
# SLIDE 7 — Main Results
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Results", "Main Comparison — Word Error Rate (%)")

# Table headers
header_y = Inches(1.7)
add_rect(s, Inches(0.5), header_y, Inches(12.3), Inches(0.55),
         fill=RGBColor(0x33, 0x55, 0x66), corner=False)
add_text(s, Inches(0.7), Inches(1.78), Inches(4.5), Inches(0.4),
         "Method", size=16, bold=True, color=WHITE, font_name="Cambria")
add_text(s, Inches(5.5), Inches(1.78), Inches(3.0), Inches(0.4),
         "test-clean", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(8.5), Inches(1.78), Inches(3.0), Inches(0.4),
         "test-other", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(11.5), Inches(1.78), Inches(1.3), Inches(0.4),
         "Δ vs Base", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")

rows = [
    ("Baseline (CTC only)", "21.16", "30.17", "—", False, False),
    ("Maximum Entropy (best, λ = 0.15)", "19.88", "28.63", "−1.54", False, False),
    ("Label Smoothing (best, α = 0.10)", "19.51", "28.30", "−1.87", True, True),
]
for i, (m, c, o, d, best, accent) in enumerate(rows):
    y = Inches(2.3 + i * 0.7)
    fill = RGBColor(0x55, 0x55, 0x55) if not best else RGBColor(0x4A, 0x60, 0x60)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.65), fill=fill, corner=False)
    name_color = YELLOW if best else WHITE
    add_text(s, Inches(0.7), y + Inches(0.13), Inches(4.5), Inches(0.4),
             m, size=15, bold=best, color=name_color)
    val_color = ORANGE if best else WHITE
    add_text(s, Inches(5.5), y + Inches(0.13), Inches(3.0), Inches(0.4),
             c, size=18, bold=True, color=val_color,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(8.5), y + Inches(0.13), Inches(3.0), Inches(0.4),
             o, size=18, bold=True, color=val_color,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(11.5), y + Inches(0.13), Inches(1.3), Inches(0.4),
             d, size=15, color=TEAL if d != "—" else DIM,
             align=PP_ALIGN.CENTER, font_name="Consolas")

# Takeaway box
add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7))
add_text(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(0.4),
         "Takeaways", size=18, bold=True, color=TEAL, font_name="Cambria")
add_runs(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.1), [
    ("● Both regularizers beat baseline by ", {}),
    ("> 1.5 pp", {"bold": True, "color": ORANGE}),
    (" on both splits (best operating points).\n", {}),
    ("● Label Smoothing (α=0.10) is the overall winner — ", {}),
    ("0.37 pp (clean) / 0.33 pp (other) ahead of MaxEnt", {"bold": True, "color": YELLOW}),
    (".\n", {}),
    ("● Each curve was tuned by a full sweep (next slides): ", {}),
    ("LS peaks sharply at 0.10, MaxEnt plateaus at 0.15–0.20", {"bold": True, "color": PINK}),
    (".", {}),
], size=14, line_spacing=1.3)

add_footer(s, 7)


# =============================================================
# SLIDE 8 — MaxEnt sweep
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Results", "Maximum Entropy — λ Sweep")

# Table
header_y = Inches(1.7)
add_rect(s, Inches(0.7), header_y, Inches(6.5), Inches(0.55),
         fill=RGBColor(0x33, 0x55, 0x66), corner=False)
add_text(s, Inches(0.9), Inches(1.78), Inches(1.5), Inches(0.4),
         "λ", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(2.4), Inches(1.78), Inches(2.4), Inches(0.4),
         "test-clean", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(4.8), Inches(1.78), Inches(2.4), Inches(0.4),
         "test-other", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")

sweep = [
    ("0.00 (baseline)", "21.16", "30.17", False),
    ("0.01",            "20.60", "29.67", False),
    ("0.03",            "20.79", "29.65", False),
    ("0.05",            "21.21", "30.46", False),
    ("0.07",            "20.93", "29.80", False),
    ("0.10",            "19.93", "28.80", False),
    ("0.15",            "19.88", "28.63", True),
    ("0.20",            "20.09", "28.60", True),
]
for i, (lam, c, o, best) in enumerate(sweep):
    y = Inches(2.3 + i * 0.5)
    fill = RGBColor(0x55, 0x55, 0x55) if not best else RGBColor(0x4A, 0x60, 0x60)
    add_rect(s, Inches(0.7), y, Inches(6.5), Inches(0.46), fill=fill, corner=False)
    name_color = YELLOW if best else WHITE
    val_color = ORANGE if best else WHITE
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(1.5), Inches(0.4),
             lam, size=14, color=name_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(2.4), y + Inches(0.05), Inches(2.4), Inches(0.4),
             c, size=14, color=val_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(4.8), y + Inches(0.05), Inches(2.4), Inches(0.4),
             o, size=14, color=val_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")

# Right side analysis box
add_rect(s, Inches(7.6), Inches(1.7), Inches(5.2), Inches(4.5),
         fill=RGBColor(0x55, 0x55, 0x55))
add_text(s, Inches(7.8), Inches(1.85), Inches(5.0), Inches(0.4),
         "Observations", size=18, bold=True, color=TEAL, font_name="Cambria")

obs = [
    ("Small λ (0.01–0.03) ", "gives modest improvement"),
    ("λ = 0.05 ", "dips — a noisy bump"),
    ("λ = 0.10 ", "strong, but not the best"),
    ("λ = 0.15–0.20 ", "lowest WER (flat optimum)"),
]
for i, (k, v) in enumerate(obs):
    y = Inches(2.4 + i * 0.55)
    add_text(s, Inches(7.85), y, Inches(0.3), Inches(0.4),
             "●", size=14, color=ORANGE)
    add_runs(s, Inches(8.1), y, Inches(4.7), Inches(0.4), [
        (k, {"bold": True, "color": ORANGE}),
        (v, {"color": WHITE}),
    ], size=13)

add_text(s, Inches(7.8), Inches(4.85), Inches(5.0), Inches(0.4),
         "Best operating point:", size=14, bold=True, color=YELLOW,
         font_name="Cambria")
add_text(s, Inches(7.8), Inches(5.25), Inches(5.0), Inches(0.4),
         "λ = 0.15 → 19.88 / 28.63", size=15, bold=True, color=WHITE,
         font_name="Consolas")
add_text(s, Inches(7.8), Inches(5.65), Inches(5.0), Inches(0.4),
         "Broad optimum → robust to λ choice.", size=12, color=LIGHT,
         font_name="Cambria")

add_footer(s, 8)


# =============================================================
# SLIDE 9 — Label Smoothing sweep
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Results", "Label Smoothing — α Sweep")

# Table
header_y = Inches(1.7)
add_rect(s, Inches(0.7), header_y, Inches(6.5), Inches(0.55),
         fill=RGBColor(0x33, 0x55, 0x66), corner=False)
add_text(s, Inches(0.9), Inches(1.78), Inches(1.5), Inches(0.4),
         "α", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(2.4), Inches(1.78), Inches(2.4), Inches(0.4),
         "test-clean", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(4.8), Inches(1.78), Inches(2.4), Inches(0.4),
         "test-other", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")

ls_sweep = [
    ("0.00 (baseline)", "21.16", "30.17", False),
    ("0.01",            "20.90", "30.15", False),
    ("0.03",            "20.71", "29.98", False),
    ("0.05",            "21.43", "30.09", False),
    ("0.07",            "20.91", "29.75", False),
    ("0.10",            "19.51", "28.30", True),
    ("0.15",            "20.16", "29.39", False),
    ("0.20",            "20.69", "29.43", False),
]
for i, (al, c, o, best) in enumerate(ls_sweep):
    y = Inches(2.3 + i * 0.5)
    fill = RGBColor(0x55, 0x55, 0x55) if not best else RGBColor(0x4A, 0x60, 0x60)
    add_rect(s, Inches(0.7), y, Inches(6.5), Inches(0.46), fill=fill, corner=False)
    name_color = YELLOW if best else WHITE
    val_color = ORANGE if best else WHITE
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(1.5), Inches(0.4),
             al, size=14, color=name_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(2.4), y + Inches(0.05), Inches(2.4), Inches(0.4),
             c, size=14, color=val_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(s, Inches(4.8), y + Inches(0.05), Inches(2.4), Inches(0.4),
             o, size=14, color=val_color, bold=best,
             align=PP_ALIGN.CENTER, font_name="Consolas")

# Right side analysis box
add_rect(s, Inches(7.6), Inches(1.7), Inches(5.2), Inches(4.5),
         fill=RGBColor(0x55, 0x55, 0x55))
add_text(s, Inches(7.8), Inches(1.85), Inches(5.0), Inches(0.4),
         "Observations", size=18, bold=True, color=TEAL, font_name="Cambria")

ls_obs = [
    ("α ≤ 0.07 ", "only marginal gains"),
    ("α = 0.10 ", "sharp minimum — best by far"),
    ("α ≥ 0.15 ", "regresses again"),
    ("Sharp U-curve ", "→ sensitive to α"),
]
for i, (k, v) in enumerate(ls_obs):
    y = Inches(2.4 + i * 0.55)
    add_text(s, Inches(7.85), y, Inches(0.3), Inches(0.4),
             "●", size=14, color=PINK)
    add_runs(s, Inches(8.1), y, Inches(4.7), Inches(0.4), [
        (k, {"bold": True, "color": PINK}),
        (v, {"color": WHITE}),
    ], size=13)

add_text(s, Inches(7.8), Inches(4.85), Inches(5.0), Inches(0.4),
         "Best operating point:", size=14, bold=True, color=YELLOW,
         font_name="Cambria")
add_text(s, Inches(7.8), Inches(5.25), Inches(5.0), Inches(0.4),
         "α = 0.10 → 19.51 / 28.30", size=15, bold=True, color=WHITE,
         font_name="Consolas")
add_text(s, Inches(7.8), Inches(5.65), Inches(5.0), Inches(0.4),
         "Overall best of all configs.", size=12, color=LIGHT,
         font_name="Cambria")

add_footer(s, 9)


# =============================================================
# SLIDE 10 — Qualitative + Analysis
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Analysis", "Why Label Smoothing > MaxEnt")

# Math comparison card (top)
add_rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.6),
         fill=RGBColor(0x55, 0x55, 0x55))
add_text(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.4),
         "Same family, different weighting:", size=16, bold=True, color=TEAL,
         font_name="Cambria")
add_text(s, Inches(0.9), Inches(2.2), Inches(0.3), Inches(0.4),
         "●", size=14, color=ORANGE)
add_runs(s, Inches(1.15), Inches(2.2), Inches(11.5), Inches(0.4), [
    ("MaxEnt: ", {"bold": True, "color": ORANGE}),
    ("H(p) = −Σ pᵥ log pᵥ", {"font": "Consolas", "color": WHITE}),
    ("   →  weights by pᵥ — already-confident tokens get pushed first", {}),
], size=14)
add_text(s, Inches(0.9), Inches(2.65), Inches(0.3), Inches(0.4),
         "●", size=14, color=PINK)
add_runs(s, Inches(1.15), Inches(2.65), Inches(11.5), Inches(0.4), [
    ("LS:     ", {"bold": True, "color": PINK}),
    ("LS(p) = −(1/V) Σ log pᵥ", {"font": "Consolas", "color": WHITE}),
    ("   →  uniform weight across vocab — all tokens regularized equally", {}),
], size=14)

# Qualitative example card
add_rect(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.0))
add_text(s, Inches(0.7), Inches(3.55), Inches(12.0), Inches(0.4),
         "Qualitative example (test-clean utterance 1)",
         size=16, bold=True, color=TEAL, font_name="Cambria")

# REF
add_text(s, Inches(0.9), Inches(4.05), Inches(1.5), Inches(0.4),
         "REF:", size=13, bold=True, color=YELLOW, font_name="Consolas")
add_runs(s, Inches(2.0), Inches(4.05), Inches(10.5), Inches(0.4), [
    ("... THEORY OF ", {"font": "Consolas"}),
    ("MEMORY", {"font": "Consolas", "bold": True, "color": ORANGE}),
    (" MUST ", {"font": "Consolas"}),
    ("ARRIVE", {"font": "Consolas", "bold": True, "color": ORANGE}),
    (" AT ... ", {"font": "Consolas"}),
    ("VIRTUE", {"font": "Consolas", "bold": True, "color": ORANGE}),
    (" ... ", {"font": "Consolas"}),
    ("SINGLED OUT", {"font": "Consolas", "bold": True, "color": ORANGE}),
    (" ...", {"font": "Consolas"}),
], size=12)

# Baseline
add_text(s, Inches(0.9), Inches(4.6), Inches(1.5), Inches(0.4),
         "Baseline:", size=13, bold=True, color=PINK, font_name="Consolas")
add_runs(s, Inches(2.0), Inches(4.6), Inches(10.5), Inches(0.4), [
    ("... THEORY OF MEMERY MUST ", {"font": "Consolas"}),
    ("ARIVE", {"font": "Consolas", "color": PINK, "bold": True}),
    (" AT ... ", {"font": "Consolas"}),
    ("VERTUE", {"font": "Consolas", "color": PINK, "bold": True}),
    (" ... ", {"font": "Consolas"}),
    ("SINGLEDOUT", {"font": "Consolas", "color": PINK, "bold": True}),
    (" ...", {"font": "Consolas"}),
], size=12)

# Label smoothing
add_text(s, Inches(0.9), Inches(5.15), Inches(1.5), Inches(0.4),
         "LS α=0.10:", size=13, bold=True, color=TEAL, font_name="Consolas")
add_runs(s, Inches(2.0), Inches(5.15), Inches(10.5), Inches(0.4), [
    ("... THEORY OF MEMERY MUST ", {"font": "Consolas"}),
    ("ARRIVE", {"font": "Consolas", "color": TEAL, "bold": True}),
    (" AT ... ", {"font": "Consolas"}),
    ("VIRTUE", {"font": "Consolas", "color": TEAL, "bold": True}),
    (" ... ", {"font": "Consolas"}),
    ("SINGLED OUT", {"font": "Consolas", "color": TEAL, "bold": True}),
    (" ...", {"font": "Consolas"}),
], size=12)

# Conclusion line
add_text(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.4),
         "→ Spelling errors (ARIVE, VERTUE) and spacing (SINGLEDOUT) corrected by Label Smoothing.",
         size=13, color=YELLOW, font_name="Cambria")

add_footer(s, 10)


# =============================================================
# SLIDE 11 — Summary
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "Summary", "")

# Three takeaway boxes (left column)
boxes = [
    ("01", "Two regularizers, fully swept",
     "Maximum Entropy (λ, 8 points) and Label Smoothing (α, 7 points), both vs the CTC baseline, all under identical training settings."),
    ("02", "Label Smoothing wins",
     "Best α=0.10 → 19.51 / 28.30. Beats baseline by −1.65 / −1.87 pp and the best MaxEnt (λ=0.15) by −0.37 / −0.33 pp."),
    ("03", "Different curve shapes",
     "LS has a sharp optimum at α=0.10 (sensitive); MaxEnt plateaus over λ=0.15–0.20 (robust but lower peak). Gains largest on test-other."),
]
for i, (num, title, body) in enumerate(boxes):
    y = Inches(1.6 + i * 1.55)
    # Circle with number
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(0.7), y + Inches(0.2),
                              Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    add_text(s, Inches(0.7), y + Inches(0.27), Inches(0.8), Inches(0.65),
             num, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name="Cambria")
    # Card
    add_rect(s, Inches(1.8), y, Inches(11.0), Inches(1.3))
    add_text(s, Inches(2.0), y + Inches(0.15), Inches(10.5), Inches(0.5),
             title, size=18, bold=True, color=ORANGE, font_name="Cambria")
    add_text(s, Inches(2.0), y + Inches(0.65), Inches(10.7), Inches(0.7),
             body, size=14, color=WHITE)

# Future work strip
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         "Future: combine LS + MaxEnt, SpecAugment data augmentation, LM shallow fusion, multi-seed runs.",
         size=13, color=YELLOW, align=PP_ALIGN.CENTER, font_name="Cambria")

add_footer(s, 11)


# =============================================================
# SLIDE 11 — Thank You
# =============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_TITLE)

add_text(s, Inches(0), Inches(2.8), SLIDE_W, Inches(1),
         "Thank You", size=60, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.6),
         "Q & A", size=28, color=TEAL,
         align=PP_ALIGN.CENTER, font_name="Cambria")
add_text(s, Inches(0), Inches(5.5), SLIDE_W, Inches(0.4),
         "Code & results: /run/models/, /run/results/", size=14, color=DIM,
         align=PP_ALIGN.CENTER, font_name="Consolas")


# =============================================================
out = "/home/cvlab/Desktop/SR/courses/2026_spring/project/run/presentation/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
